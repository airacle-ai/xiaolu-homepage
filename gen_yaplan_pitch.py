#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成 yaplan-pitch.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
import copy
from lxml import etree

# ─────────────────────────────────────────
# 路径
# ─────────────────────────────────────────
TEMPLATE_PATH = "/home/aitist/.omnara/attachments/68fb5f1e-0cb1-4990-8dcc-090bff582919/AI______PPT__-a1c0b1595097a4d7c911b423.pptx"
OUTPUT_PATH   = "/mnt/c/Users/oyzh8/Desktop/win_code/docker/omnara_workspace/steve/xiaolu-homepage/yaplan-pitch.pptx"

# ─────────────────────────────────────────
# 颜色常量
# ─────────────────────────────────────────
C_DARK_BLUE  = RGBColor(0x08, 0x25, 0x4D)
C_RED        = RGBColor(0xC0, 0x00, 0x00)
C_BODY       = RGBColor(0x33, 0x33, 0x33)
C_GRAY       = RGBColor(0x66, 0x66, 0x66)
C_ORANGE     = RGBColor(0xFF, 0x8C, 0x42)
C_GREEN      = RGBColor(0x4E, 0xBF, 0x8F)
C_BLUE       = RGBColor(0x4A, 0x90, 0xD9)
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK      = RGBColor(0x00, 0x00, 0x00)
C_LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
C_DARK_GRAY  = RGBColor(0x99, 0x99, 0x99)
C_LIGHT_BLUE_BG = RGBColor(0xF7, 0xFA, 0xFF)
C_LIGHT_ORG_BG  = RGBColor(0xFF, 0xF8, 0xF5)
C_LIGHT_GRN_BG  = RGBColor(0xF5, 0xFF, 0xF8)

TOTAL_SLIDES = 12
FONT_NAME    = "微软雅黑"

# ─────────────────────────────────────────
# 辅助函数
# ─────────────────────────────────────────

def set_font(run_or_para, font_name=FONT_NAME, size=None, bold=None,
             color=None, italic=None):
    """统一设置字体，支持 run 或 paragraph"""
    if hasattr(run_or_para, 'font'):
        font = run_or_para.font
    else:
        return
    font.name = font_name
    # 中文字体需要同时设置 eastAsian
    try:
        rPr = font._element
        rPr.set(qn('w:eastAsia'), font_name)
    except Exception:
        pass
    if size is not None:
        font.size = Pt(size)
    if bold is not None:
        font.bold = bold
    if color is not None:
        font.color.rgb = color
    if italic is not None:
        font.italic = italic


def _set_para_font(para, font_name=FONT_NAME, size=None, bold=None,
                   color=None, align=None, italic=None):
    if align is not None:
        if align == 'center':
            para.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            para.alignment = PP_ALIGN.RIGHT
        else:
            para.alignment = PP_ALIGN.LEFT
    for run in para.runs:
        set_font(run, font_name=font_name, size=size, bold=bold,
                 color=color, italic=italic)


def add_rect(slide, left, top, width, height,
             fill_color=None, line_color=None, line_width=None,
             round_corners=False):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    fill = shape.fill
    if fill_color is not None:
        fill.solid()
        fill.fore_color.rgb = fill_color
    else:
        fill.background()

    line = shape.line
    if line_color is not None:
        line.color.rgb = line_color
        if line_width is not None:
            line.width = Pt(line_width)
    else:
        line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=None,
                align='left', wrap=True, italic=False,
                line_spacing=None):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap

    # 清空默认段落
    p = tf.paragraphs[0]
    p.text = text

    if align == 'center':
        p.alignment = PP_ALIGN.CENTER
    elif align == 'right':
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    for run in p.runs:
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
        run.font.italic = italic

    return txBox


def add_multiline_textbox(slide, lines, left, top, width, height,
                          font_size=12, bold=False, color=None,
                          align='left', line_spacing_pt=None):
    """添加多行文本框（lines是列表）"""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        if align == 'center':
            p.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        for run in p.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color

    return txBox


def add_title(slide, title_text, subtitle=None, slide_num=1,
              tag_text=None):
    """添加标题+左侧竖线+右下页码+品牌"""
    # 左侧竖线
    line_shape = slide.shapes.add_shape(
        1,
        Inches(0.47), Inches(0.35),
        Inches(0.05), Inches(6.5)
    )
    line_shape.fill.solid()
    line_shape.fill.fore_color.rgb = C_RED
    line_shape.line.fill.background()

    # 标题
    title_box = add_textbox(
        slide, title_text,
        left=0.65, top=0.35, width=9.5, height=0.6,
        font_size=24, bold=True, color=C_DARK_BLUE, align='left'
    )

    if subtitle:
        add_textbox(
            slide, subtitle,
            left=0.65, top=0.92, width=9.5, height=0.35,
            font_size=12, bold=False, color=C_GRAY, align='left'
        )

    # 右下页码
    num_str = f"{slide_num:02d} / {TOTAL_SLIDES:02d}"
    add_textbox(
        slide, num_str,
        left=11.8, top=7.1, width=1.3, height=0.3,
        font_size=9, bold=False, color=C_DARK_GRAY, align='right'
    )

    # 右下品牌
    add_textbox(
        slide, "芽计 YaPlan",
        left=10.2, top=7.1, width=1.5, height=0.3,
        font_size=9, bold=False, color=C_DARK_GRAY, align='right'
    )

    # 章节标签（右上角）
    if tag_text:
        tag_bg = add_rect(
            slide, left=10.5, top=0.15, width=2.6, height=0.35,
            fill_color=C_DARK_GRAY, line_color=None
        )
        add_textbox(
            slide, tag_text,
            left=10.5, top=0.15, width=2.6, height=0.35,
            font_size=9, bold=False, color=C_WHITE, align='center'
        )


def add_hline(slide, left, top, width, color=C_LIGHT_GRAY, height=0.02):
    """添加水平分隔线"""
    shape = add_rect(slide, left, top, width, height,
                     fill_color=color, line_color=None)
    return shape


# ─────────────────────────────────────────
# 打开模板，创建演示文稿
# ─────────────────────────────────────────
print("正在加载模板...")
tpl = Presentation(TEMPLATE_PATH)
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# 使用空白布局
blank_layout = prs.slide_layouts[6]  # 通常index 6是空白

# ═══════════════════════════════════════════════════════════════
# Slide 1 — 产品画布（Lean Canvas）
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 1 — 产品画布...")
slide = prs.slides.add_slide(blank_layout)

add_title(slide, "产品画布 · Lean Canvas", slide_num=1, tag_text="产品画布")

# 三列大矩形
# 左列：蓝
add_rect(slide, left=0.47, top=1.4, width=3.9, height=5.4,
         fill_color=C_LIGHT_BLUE_BG,
         line_color=C_BLUE, line_width=1.5)
# 中列：橙
add_rect(slide, left=4.52, top=1.4, width=3.9, height=5.4,
         fill_color=C_LIGHT_ORG_BG,
         line_color=C_ORANGE, line_width=1.5)
# 右列：绿
add_rect(slide, left=8.57, top=1.4, width=3.9, height=5.4,
         fill_color=C_LIGHT_GRN_BG,
         line_color=C_GREEN, line_width=1.5)

# 左列内容
add_textbox(slide, "🎯 问题与痛点",
            left=0.67, top=1.5, width=3.5, height=0.45,
            font_size=14, bold=True, color=C_BLUE)
add_hline(slide, left=0.67, top=1.95, width=3.5)
add_multiline_textbox(slide,
    ["① 教育支出焦虑，缺乏整体规划框架",
     "② 保险配置混乱，面对销售话术无从判断",
     "③ 焦虑型课外消费泛滥，38%支出属情绪驱动",
     "④ 零花钱给了不知怎么引导"],
    left=0.67, top=2.0, width=3.5, height=2.8,
    font_size=12, color=C_BODY)

# 中列内容
add_textbox(slide, "💡 解决方案（攒·保·理·教）",
            left=4.72, top=1.5, width=3.5, height=0.45,
            font_size=13, bold=True, color=C_ORANGE)
add_hline(slide, left=4.72, top=1.95, width=3.5)
add_multiline_textbox(slide,
    ["攒：AI双情景教育金测算与定投计划",
     "保：保单OCR解读，非销售保险配置建议",
     "理：多模态记账+焦虑消费3层检测+冷静期",
     "教：三罐子零花钱+劳动奖励+财商小课堂"],
    left=4.72, top=2.0, width=3.5, height=2.8,
    font_size=12, color=C_BODY)

# 右列内容
add_textbox(slide, "✦ 独特价值主张",
            left=8.77, top=1.5, width=3.5, height=0.45,
            font_size=14, bold=True, color=C_GREEN)
add_hline(slide, left=8.77, top=1.95, width=3.5)
add_multiline_textbox(slide,
    ['"第一个对你说\'不\'的财务顾问"',
     "",
     "纯规划、不销售",
     "焦虑消费识别（行业首创）",
     "生活化比喻翻译金融术语",
     "0-18岁全周期陪伴"],
    left=8.77, top=2.0, width=3.5, height=2.8,
    font_size=12, color=C_BODY)

# 底部三个信息块
for bx, bt, bc in [
    (0.47,  "用户群体：0-18岁子女城市中产家庭", C_BLUE),
    (4.52,  "商业模式：Freemium ¥99-199/月",     C_ORANGE),
    (8.57,  "技术：Claude API + 多模态OCR + 规则引擎", C_GREEN),
]:
    add_rect(slide, left=bx, top=6.55, width=3.9, height=0.7,
             fill_color=RGBColor(0xF0, 0xF0, 0xF0), line_color=bc, line_width=1)
    add_textbox(slide, bt,
                left=bx+0.1, top=6.6, width=3.7, height=0.55,
                font_size=10, color=C_BODY, align='center')

print("  Slide 1 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 2 — 用户需求·三阶段画像
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 2 — 三阶段家庭画像...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "用户需求分析 · 三阶段家庭画像",
          slide_num=2, tag_text="用户研究 01/03")

col_data = [
    (0.47,  C_BLUE,   "新生期 0-3岁",   "👤 25-35岁双职工，月入1.5-4万",
     ["• 早教课值不值？报还是不报？",
      "• 保险该买哪些？被话术绑架",
      "• 教育金何时开始？怎么存？"], "¥3,200"),
    (4.43,  C_ORANGE, "成长期 4-12岁",  "👤 30-42岁中产，月入2.5-8万",
     ["• 课外班军备竞赛，月花8,000+",
      "• 教育金储备进度滞后",
      "• 同伴压力导致跟风报班"], "¥6,800"),
    (8.40,  C_GREEN,  "冲刺期 13-18岁", "👤 40-52岁，月入4-12万",
     ["• 出国vs国内？200-500万差距",
      "• 留学费用测算，缺乏独立工具",
      "• 资产不动产化，流动性不足"], "¥12,400"),
]

for lx, color, period, user_desc, pain_points, spend in col_data:
    # 顶部色块
    add_rect(slide, left=lx, top=1.5, width=3.8, height=0.55,
             fill_color=color)
    add_textbox(slide, period,
                left=lx+0.1, top=1.52, width=3.6, height=0.48,
                font_size=15, bold=True, color=C_WHITE, align='center')
    # 内容区
    add_rect(slide, left=lx, top=2.05, width=3.8, height=4.45,
             fill_color=C_WHITE,
             line_color=color, line_width=1)
    add_textbox(slide, user_desc,
                left=lx+0.15, top=2.1, width=3.5, height=0.4,
                font_size=10, color=C_GRAY)
    add_hline(slide, left=lx+0.1, top=2.55, width=3.6, color=C_LIGHT_GRAY)
    add_textbox(slide, "核心痛点：",
                left=lx+0.15, top=2.65, width=3.5, height=0.3,
                font_size=11, bold=True, color=C_DARK_BLUE)
    add_multiline_textbox(slide, pain_points,
                          left=lx+0.15, top=3.0, width=3.5, height=1.8,
                          font_size=11, color=C_BODY)
    add_hline(slide, left=lx+0.1, top=4.85, width=3.6, color=C_LIGHT_GRAY)
    add_textbox(slide, "月均教育支出",
                left=lx+0.15, top=4.92, width=3.5, height=0.3,
                font_size=10, color=C_GRAY)
    add_textbox(slide, spend,
                left=lx+0.15, top=5.2, width=3.5, height=0.6,
                font_size=22, bold=True, color=C_RED)

print("  Slide 2 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 3 — 用户需求·焦虑消费识别
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 3 — 焦虑消费识别...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "用户需求分析 · 焦虑型消费识别",
          slide_num=3, tag_text="用户研究 02/03")

# 顶部数据横幅（3个数字卡片）
stats = [
    ("38%",   C_RED,    "育儿支出属焦虑驱动"),
    ("87%",   C_ORANGE, "家长报名后曾后悔的课程"),
    ("-¥2,100", C_GREEN, "芽计用户月均识别的焦虑消费"),
]
for i, (num, color, desc) in enumerate(stats):
    lx = 0.47 + i * 4.3
    add_rect(slide, left=lx, top=1.25, width=4.1, height=1.1,
             fill_color=C_LIGHT_GRAY, line_color=color, line_width=1.5)
    add_textbox(slide, num,
                left=lx+0.1, top=1.28, width=2.0, height=0.65,
                font_size=28, bold=True, color=color)
    add_textbox(slide, desc,
                left=lx+0.1, top=1.9, width=3.8, height=0.35,
                font_size=11, color=C_GRAY)

# 左侧红色卡片
add_rect(slide, left=0.47, top=2.55, width=5.5, height=3.5,
         fill_color=RGBColor(0xFF, 0xF5, 0xF5),
         line_color=None)
# 左边框
add_rect(slide, left=0.47, top=2.55, width=0.08, height=3.5,
         fill_color=C_RED)
add_textbox(slide, "❌ 焦虑驱动型消费",
            left=0.65, top=2.65, width=5.0, height=0.45,
            font_size=14, bold=True, color=C_RED)
add_multiline_textbox(slide,
    ["• 看到KOL推荐冲动下单",
     "• 听说别人家孩子都在学",
     "• 周末深夜突击报班",
     "• 无ROI评估，跟风即买",
     "• 孩子不感兴趣但续费"],
    left=0.65, top=3.15, width=5.1, height=2.5,
    font_size=12, color=C_BODY)

# 中间"芽计帮你识别"标签
add_rect(slide, left=5.9, top=2.55, width=1.4, height=3.5,
         fill_color=C_ORANGE)
add_textbox(slide, "芽\n计\n帮\n你\n识\n别",
            left=5.95, top=2.8, width=1.3, height=3.0,
            font_size=13, bold=True, color=C_WHITE, align='center')

# 右侧绿色卡片
add_rect(slide, left=7.4, top=2.55, width=5.5, height=3.5,
         fill_color=RGBColor(0xF5, 0xFF, 0xF8),
         line_color=None)
add_rect(slide, left=7.4, top=2.55, width=0.08, height=3.5,
         fill_color=C_GREEN)
add_textbox(slide, "✅ 理性规划型消费",
            left=7.6, top=2.65, width=5.0, height=0.45,
            font_size=14, bold=True, color=C_GREEN)
add_multiline_textbox(slide,
    ["• 基于孩子兴趣和意愿决策",
     "• 有明确学习目标和退出机制",
     "• 评估性价比后决定",
     "• 定期复盘课程出勤与效果",
     "• 纳入家庭月度预算管控"],
    left=7.6, top=3.15, width=5.1, height=2.5,
    font_size=12, color=C_BODY)

# 底部说明
add_textbox(slide,
            "芽计通过行为信号+语义情绪+综合评分三层机制，在消费决策前给出理性提醒",
            left=0.7, top=6.3, width=12.0, height=0.4,
            font_size=10, color=C_DARK_GRAY, align='center')

print("  Slide 3 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 4 — 用户需求·市场规模
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 4 — 市场规模...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "用户需求分析 · 市场规模与机会",
          slide_num=4, tag_text="用户研究 03/03")

# 左侧：三层嵌套矩形 TAM/SAM/SOM
tam_data = [
    (0.47, 1.35, 8.0, 5.3, RGBColor(0xEE, 0xF3, 0xFF),
     "TAM  总市场：2.1亿有子女家庭 · 潜在规模¥6,300亿",
     C_BLUE),
    (1.0,  1.9,  7.0, 4.1, RGBColor(0xFF, 0xF3, 0xE8),
     "SAM  可触达市场：4,200万城市中产家庭 · ¥1,260亿",
     C_ORANGE),
    (1.7,  2.55, 5.6, 2.8, C_ORANGE,
     "SOM  目标：210万家庭 · ¥6.3亿ARR（3年）",
     C_WHITE),
]
for lx, tp, w, h, bg, text, tc in tam_data:
    add_rect(slide, left=lx, top=tp, width=w, height=h,
             fill_color=bg, line_color=None)
    add_textbox(slide, text,
                left=lx+0.2, top=tp+0.1, width=w-0.3, height=h-0.15,
                font_size=12, bold=(tc == C_WHITE), color=tc)

# 右侧：4个数据卡片
right_cards = [
    ("68%",    C_RED,    "受访家长表示缺乏系统财务规划"),
    ("¥3.2万", C_BLUE,   "城市家庭年均课外教育支出"),
    ("92%",    C_GREEN,  "愿为精准规划付费的家长比例"),
    ("18年",   C_ORANGE, "单用户生命周期 · LTV可达¥3,582"),
]
for i, (num, color, desc) in enumerate(right_cards):
    tp = 1.35 + i * 1.25
    add_rect(slide, left=8.8, top=tp, width=4.2, height=1.1,
             fill_color=C_WHITE, line_color=color, line_width=1.5)
    add_rect(slide, left=8.8, top=tp, width=0.08, height=1.1,
             fill_color=color)
    add_textbox(slide, num,
                left=9.0, top=tp+0.05, width=1.8, height=0.65,
                font_size=22, bold=True, color=color)
    add_textbox(slide, desc,
                left=9.0, top=tp+0.65, width=3.8, height=0.35,
                font_size=10, color=C_GRAY)

# 底部数据来源
add_textbox(slide,
            "数据来源：艾瑞咨询2025 / 国家统计局 / CFPS家庭追踪调查",
            left=0.6, top=6.9, width=9.0, height=0.35,
            font_size=9, color=C_DARK_GRAY)

print("  Slide 4 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 5 — 产品定位与差异化
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 5 — 产品定位与差异化...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "产品方案 · 攒·保·理·教 四模块定位",
          slide_num=5, tag_text="产品方案 01/07")

# 顶部框架横幅
add_rect(slide, left=0.47, top=1.3, width=12.5, height=0.5,
         fill_color=C_LIGHT_GRAY)
add_textbox(slide,
            "所有功能通过对话式AI串联，构成亲子家庭财务的整体框架 · 角色定位：家庭财务医生 × 成长陪伴学姐",
            left=0.6, top=1.32, width=12.2, height=0.45,
            font_size=11, color=C_BODY, align='center')

# 四个模块卡片（2×2）
modules = [
    (0.47, 1.95, C_BLUE,   "📈 攒 · 教育金规划",
     "PMT公式精准测算缺口 / 双情景模拟（国内vs出国）\n每年开学季动态校准 / 生活化比喻：¥36万缺口=360次家庭大餐"),
    (6.3,  1.95, C_RED,    "🛡️ 保 · 保险配置",
     "保单OCR解读关键条款 / 先大人后小孩优先级\n非销售承诺，只分析不导流 / 通俗险种比喻库"),
    (0.47, 4.3,  C_ORANGE, "💸 理 · 支出管理",
     "拍票据/截图/语音多模态记账 / 3层焦虑消费检测\n24h购物车冷静期 / 兴趣班ROI热力表"),
    (6.3,  4.3,  C_GREEN,  "💰 教 · 财商教育",
     "三罐子零花钱模型 / 劳动-奖励任务清单\nAI发零花钱时推财商故事 / 财商绘本讲给AI听"),
]

for lx, tp, border_color, title, content in modules:
    add_rect(slide, left=lx, top=tp, width=5.7, height=2.2,
             fill_color=C_WHITE, line_color=border_color, line_width=2)
    add_textbox(slide, title,
                left=lx+0.2, top=tp+0.1, width=5.3, height=0.45,
                font_size=14, bold=True, color=border_color)
    add_hline(slide, left=lx+0.15, top=tp+0.6, width=5.4,
              color=border_color, height=0.02)
    add_textbox(slide, content,
                left=lx+0.2, top=tp+0.65, width=5.3, height=1.4,
                font_size=11, color=C_BODY)

print("  Slide 5 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 6 — 竞品对比与差异化
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 6 — 竞品对比...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "产品方案 · 差异化定位",
          slide_num=6, tag_text="产品方案 02/07")

# 竞品对比表
table_data = [
    ["功能维度",       "芽计 YaPlan",    "挖财",      "支付宝财富",  "Greenlight"],
    ["焦虑消费识别",   "✅ 行业首创",    "❌",         "❌",          "❌"],
    ["保单OCR解读",    "✅ 通俗解析",    "❌",         "❌",          "❌"],
    ["非销售导向",     "✅ 无返佣",      "❌ 有广告",  "❌ 有佣金",   "✅ 基础"],
    ["教育金测算",     "✅ 双情景",      "无",         "弱",          "无"],
    ["亲子财商教育",   "✅ 三罐子+绘本", "❌",         "❌",          "✅ 基础"],
    ["生活化比喻",     "✅ 内置系统",    "❌",         "❌",          "❌"],
]

rows = len(table_data)
cols = len(table_data[0])
col_widths = [2.5, 2.8, 1.8, 1.8, 1.8]  # inches
table_top  = 1.35
table_left = 0.47
row_height = 0.55

tbl = slide.shapes.add_table(
    rows, cols,
    Inches(table_left), Inches(table_top),
    Inches(sum(col_widths)), Inches(rows * row_height)
).table

# 设置列宽
for ci, cw in enumerate(col_widths):
    tbl.columns[ci].width = Inches(cw)

for ri, row in enumerate(table_data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.text = val
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        for run in tf.paragraphs[0].runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(11)
            if ri == 0:
                run.font.bold = True
                run.font.color.rgb = C_WHITE
            elif ci == 1:
                run.font.bold = True
                if "✅" in val:
                    run.font.color.rgb = C_GREEN
                elif "❌" in val:
                    run.font.color.rgb = C_DARK_GRAY
                else:
                    run.font.color.rgb = C_BODY
            elif "✅" in val:
                run.font.color.rgb = C_GREEN
            elif "❌" in val:
                run.font.color.rgb = C_DARK_GRAY
            else:
                run.font.color.rgb = C_BODY

        # 背景色
        fill = cell.fill
        fill.solid()
        if ri == 0:
            fill.fore_color.rgb = C_DARK_BLUE
        elif ci == 1:
            fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xE8)
        elif ri % 2 == 0:
            fill.fore_color.rgb = C_LIGHT_GRAY
        else:
            fill.fore_color.rgb = C_WHITE

# 底部三个差异化标签
diff_tags = [
    (C_BLUE,   "🧠 焦虑识别引擎",   "行业首创"),
    (C_ORANGE, "📊 非销售导向",      "信任经济壁垒"),
    (C_GREEN,  "💬 语言翻译器",      "金融术语→生活比喻"),
]
for i, (color, title, sub) in enumerate(diff_tags):
    lx = 0.47 + i * 4.3
    add_rect(slide, left=lx, top=5.75, width=3.9, height=1.0,
             fill_color=color, line_color=None)
    add_textbox(slide, title,
                left=lx+0.1, top=5.8, width=3.7, height=0.45,
                font_size=13, bold=True, color=C_WHITE, align='center')
    add_textbox(slide, sub,
                left=lx+0.1, top=6.2, width=3.7, height=0.4,
                font_size=11, color=C_WHITE, align='center')

print("  Slide 6 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 7 — AI能力运用
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 7 — AI能力运用...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "产品方案 · AI能力运用",
          slide_num=7, tag_text="产品方案 03/07")

# 顶部AI横幅
add_rect(slide, left=0.47, top=1.3, width=12.5, height=0.5,
         fill_color=C_LIGHT_GRAY)
add_textbox(slide,
            "AI双重角色：🏥 家庭财务医生（客观、专业、不卖东西） × 👩‍🎓 成长陪伴学姐（温暖、真实、懂孩子）",
            left=0.6, top=1.32, width=12.2, height=0.45,
            font_size=11, color=C_BODY, align='center')

# 左侧：5步流程
steps = [
    ("Step1", "🗣️ 对话采集", "自然语言问答，采集家庭基本信息"),
    ("Step2", "🔍 需求挖掘", "识别显性需求与隐性焦虑来源"),
    ("Step3", "📊 方案生成", "PMT计算+双情景+生活化比喻输出"),
    ("Step4", "⚠️ 焦虑检测", "行为信号+语义情绪+综合评分三层"),
    ("Step5", "🔄 动态复盘", "月度小报+季度报告+年度校准"),
]
step_colors = [C_BLUE, C_BLUE, C_ORANGE, C_RED, C_GREEN]

for i, ((sn, title, desc), color) in enumerate(zip(steps, step_colors)):
    tp = 1.95 + i * 0.95
    add_rect(slide, left=0.47, top=tp, width=6.5, height=0.8,
             fill_color=color, line_color=None)
    add_textbox(slide, f"{sn}  {title}",
                left=0.6, top=tp+0.03, width=3.0, height=0.4,
                font_size=12, bold=True, color=C_WHITE)
    add_textbox(slide, desc,
                left=0.6, top=tp+0.4, width=6.1, height=0.35,
                font_size=11, color=C_WHITE)
    # 向下箭头（除最后一步）
    if i < 4:
        add_textbox(slide, "▼",
                    left=3.4, top=tp+0.8, width=0.5, height=0.2,
                    font_size=10, color=C_GRAY, align='center')

# 右侧：多模态能力卡片
mm_caps = [
    ("📸 拍票据", "OCR自动归类"),
    ("🎙️ 语音记账", "自然语言录入"),
    ("📄 截图保单", "AI提取条款"),
]
for i, (title, sub) in enumerate(mm_caps):
    lx = 7.3 + i * 1.95
    add_rect(slide, left=lx, top=1.95, width=1.8, height=1.5,
             fill_color=RGBColor(0xF0, 0xF4, 0xFF),
             line_color=C_BLUE, line_width=1)
    add_textbox(slide, title,
                left=lx+0.05, top=2.05, width=1.7, height=0.5,
                font_size=12, bold=True, color=C_DARK_BLUE, align='center')
    add_textbox(slide, sub,
                left=lx+0.05, top=2.55, width=1.7, height=0.5,
                font_size=10, color=C_GRAY, align='center')

# 底部对话示例
add_rect(slide, left=7.3, top=3.6, width=5.8, height=1.2,
         fill_color=RGBColor(0xEE, 0xF3, 0xFF),
         line_color=C_BLUE, line_width=1)
add_textbox(slide,
            "用户：'5岁女儿，想出国，月入3-5万'\n→ 芽计：'目标¥120万，月投¥4,900，\n相当于每天少一杯奶茶+咖啡'",
            left=7.4, top=3.65, width=5.6, height=1.1,
            font_size=11, color=C_DARK_BLUE)

# 右侧底部技术标注
add_textbox(slide, "底层：Claude API · 多模态OCR · 规则引擎",
            left=7.3, top=5.0, width=5.8, height=0.4,
            font_size=10, color=C_GRAY)

print("  Slide 7 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 8 — 商业模式
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 8 — 商业模式...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "产品方案 · 商业模式",
          slide_num=8, tag_text="产品方案 04/07")

# 三列定价卡片
pricing = [
    (0.5,  C_GRAY,         "免费版 Free",    "",
     ["基础教育金测算", "AI对话3次/月", "支出分类记录", "焦虑消费基础识别"],
     False),
    (4.5,  C_ORANGE,       "标准版 ¥99/月",  "推荐",
     ["无限AI规划对话", "焦虑消费深度识别", "保单OCR解读",
      "月度亲子财务小报", "兴趣班ROI热力表"],
     True),
    (8.5,  C_DARK_BLUE,    "高级版 ¥199/月", "",
     ["全部标准版功能", "专属1v1顾问", "家庭多账号",
      "年度财务体检报告", "优先体验新功能"],
     False),
]

for lx, color, plan, badge, features, highlight in pricing:
    tp = 1.4 if not highlight else 1.25
    h  = 4.0 if not highlight else 4.3
    add_rect(slide, left=lx, top=tp, width=3.8, height=h,
             fill_color=C_WHITE, line_color=color, line_width=2)
    # 顶部色块
    add_rect(slide, left=lx, top=tp, width=3.8, height=0.6,
             fill_color=color)
    add_textbox(slide, plan,
                left=lx+0.1, top=tp+0.05, width=3.5, height=0.5,
                font_size=14, bold=True, color=C_WHITE, align='center')
    # 推荐角标
    if badge:
        add_rect(slide, left=lx+2.9, top=tp, width=0.9, height=0.35,
                 fill_color=C_RED)
        add_textbox(slide, badge,
                    left=lx+2.9, top=tp+0.02, width=0.9, height=0.3,
                    font_size=9, bold=True, color=C_WHITE, align='center')
    # 功能列表
    for j, feat in enumerate(features):
        add_textbox(slide, f"✓  {feat}",
                    left=lx+0.2, top=tp+0.75+j*0.55, width=3.4, height=0.45,
                    font_size=11, color=C_BODY)

# 底部增长飞轮
flywheel = ["免费体验", "价值验证", "付费转化", "口碑传播", "更多用户"]
for i, node in enumerate(flywheel):
    lx = 0.7 + i * 2.4
    add_rect(slide, left=lx, top=5.7, width=1.9, height=0.7,
             fill_color=C_ORANGE, line_color=None)
    add_textbox(slide, node,
                left=lx, top=5.75, width=1.9, height=0.6,
                font_size=11, bold=True, color=C_WHITE, align='center')
    if i < 4:
        add_textbox(slide, "→",
                    left=lx+1.9, top=5.85, width=0.4, height=0.4,
                    font_size=14, bold=True, color=C_ORANGE, align='center')

print("  Slide 8 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 9 — 原型展示①（首页+AI）
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 9 — 原型展示①...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "产品原型 · 首页 & AI对话",
          slide_num=9, tag_text="产品方案 05/07")


def draw_phone(slide, lx, tp, screen_title, ui_items):
    """绘制手机轮廓+模拟UI"""
    pw, ph = 2.8, 5.0
    # 手机外框
    add_rect(slide, left=lx, top=tp, width=pw, height=ph,
             fill_color=RGBColor(0x1A, 0x1A, 0x1A), line_color=None)
    # 屏幕区域
    add_rect(slide, left=lx+0.15, top=tp+0.25, width=pw-0.3, height=ph-0.5,
             fill_color=C_WHITE, line_color=None)
    # 顶部状态栏
    add_rect(slide, left=lx+0.15, top=tp+0.25, width=pw-0.3, height=0.25,
             fill_color=C_DARK_BLUE, line_color=None)
    add_textbox(slide, screen_title,
                left=lx+0.15, top=tp+0.27, width=pw-0.3, height=0.22,
                font_size=7, bold=True, color=C_WHITE, align='center')
    # UI元素
    for j, (item_text, item_color) in enumerate(ui_items):
        iy = tp + 0.6 + j * 0.72
        add_rect(slide, left=lx+0.2, top=iy, width=pw-0.4, height=0.6,
                 fill_color=item_color, line_color=None)
        add_textbox(slide, item_text,
                    left=lx+0.25, top=iy+0.08, width=pw-0.5, height=0.45,
                    font_size=8, color=C_BODY)


# 左侧手机：首页Dashboard
draw_phone(slide, lx=0.7, tp=1.5,
           screen_title="🏠 首页 · Dashboard",
           ui_items=[
               ("教育金健康度 18%  ⚠️",           RGBColor(0xFF, 0xF5, 0xF5)),
               ("本月支出 ¥8,200  焦虑消费31%",   RGBColor(0xFF, 0xF8, 0xF0)),
               ("教育金缺口 ¥36万  立即规划 →",   RGBColor(0xF0, 0xF8, 0xFF)),
               ("快速入口：记账 / 规划 / 保险",    RGBColor(0xF5, 0xFF, 0xF8)),
               ("宝贝 · 本周财商故事",             RGBColor(0xFF, 0xF8, 0xFF)),
           ])
add_textbox(slide, "🏠 首页 · Dashboard",
            left=0.7, top=6.55, width=2.8, height=0.35,
            font_size=10, color=C_DARK_BLUE, align='center')

# 右侧手机：AI对话
draw_phone(slide, lx=5.5, tp=1.5,
           screen_title="🤖 AI · 对话规划",
           ui_items=[
               ("你好！我来帮你规划教育金～",       RGBColor(0xF0, 0xF4, 0xFF)),
               ("孩子几岁？目标出国还是国内？",     RGBColor(0xF0, 0xF4, 0xFF)),
               ("5岁女儿，希望出国留学",            RGBColor(0xFF, 0xFF, 0xFF)),
               ("目标¥120万，月投¥4,900起",        RGBColor(0xF5, 0xFF, 0xF8)),
               ("相当于每天少一杯奶茶+咖啡 ☕",    RGBColor(0xF5, 0xFF, 0xF8)),
           ])
add_textbox(slide, "🤖 AI · 对话规划",
            left=5.5, top=6.55, width=2.8, height=0.35,
            font_size=10, color=C_DARK_BLUE, align='center')

# 右侧标注区
highlights = [
    "教育金健康度18%，实时预警",
    "焦虑消费占比31%，红色标注",
    "7轮对话生成双情景规划报告",
    "生活化比喻：缺口=360次家庭大餐",
]
for i, h in enumerate(highlights):
    tp = 1.8 + i * 1.1
    add_rect(slide, left=9.0, top=tp, width=0.2, height=0.2,
             fill_color=C_ORANGE)
    add_textbox(slide, h,
                left=9.3, top=tp-0.05, width=3.8, height=0.65,
                font_size=11, color=C_BODY)

print("  Slide 9 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 10 — 原型展示②（账单+规划）
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 10 — 原型展示②...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "产品原型 · 账单分析 & 教育金规划",
          slide_num=10, tag_text="产品方案 06/07")

draw_phone(slide, lx=0.7, tp=1.5,
           screen_title="💳 账单 · 支出分析",
           ui_items=[
               ("课外班 41% ⚠️ 超均值13%",        RGBColor(0xFF, 0xF5, 0xF5)),
               ("生活必需 28%  餐饮 15%",          RGBColor(0xF5, 0xFF, 0xF8)),
               ("焦虑消费预警：奥数班 ⚠️",          RGBColor(0xFF, 0xF8, 0xF0)),
               ("焦虑消费预警：英语外教 ⚠️",        RGBColor(0xFF, 0xF8, 0xF0)),
               ("本月节省焦虑支出 ¥2,100",         RGBColor(0xF0, 0xF8, 0xFF)),
           ])
add_textbox(slide, "💳 账单 · 支出分析",
            left=0.7, top=6.55, width=2.8, height=0.35,
            font_size=10, color=C_DARK_BLUE, align='center')

draw_phone(slide, lx=5.5, tp=1.5,
           screen_title="📋 规划 · 教育金+保险",
           ui_items=[
               ("教育金  |  保险配置",              RGBColor(0xF0, 0xF4, 0xFF)),
               ("教育金目标 ¥120万  完成48%",       RGBColor(0xF5, 0xFF, 0xF8)),
               ("建议月投增至 ¥5,200 (+¥300)",     RGBColor(0xFF, 0xF8, 0xF0)),
               ("上传保单 → AI解读中…",             RGBColor(0xF5, 0xF5, 0xFF)),
               ("重疾保额不足，建议补充¥30万",      RGBColor(0xFF, 0xF5, 0xF5)),
           ])
add_textbox(slide, "📋 规划 · 教育金+保险",
            left=5.5, top=6.55, width=2.8, height=0.35,
            font_size=10, color=C_DARK_BLUE, align='center')

highlights2 = [
    "五色饼图：课外班占41%，超均值13%",
    "焦虑识别：奥数班+英语外教双预警",
    "教育金/保险二级Tab切换",
    "保单OCR上传，AI解读重疾保额不足",
]
for i, h in enumerate(highlights2):
    tp = 1.8 + i * 1.1
    add_rect(slide, left=9.0, top=tp, width=0.2, height=0.2,
             fill_color=C_ORANGE)
    add_textbox(slide, h,
                left=9.3, top=tp-0.05, width=3.8, height=0.65,
                font_size=11, color=C_BODY)

print("  Slide 10 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 11 — 原型展示③（财商教育）
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 11 — 原型展示③...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "产品原型 · 财商教育模块",
          slide_num=11, tag_text="产品方案 07/07")

draw_phone(slide, lx=0.8, tp=1.5,
           screen_title="💰 财商教育",
           ui_items=[
               ("🫙 三罐子零花钱",                 RGBColor(0xFF, 0xF8, 0xF0)),
               ("自由¥12 储蓄¥6 分享¥2",          RGBColor(0xFF, 0xF8, 0xF0)),
               ("✅ 今日任务：整理房间 +¥5",        RGBColor(0xF5, 0xFF, 0xF8)),
               ("📖 今日财商故事：懒蚂蚁",          RGBColor(0xF0, 0xF4, 0xFF)),
               ("宝贝成长报告 · 本月 ⭐⭐⭐",       RGBColor(0xFF, 0xFF, 0xF0)),
           ])
add_textbox(slide, "💰 财商教育",
            left=0.8, top=6.55, width=2.8, height=0.35,
            font_size=10, color=C_DARK_BLUE, align='center')

# 右侧三个功能卡片
fn_cards = [
    (C_ORANGE, "🫙 三罐子零花钱",
     "自由支配罐¥12 / 储蓄罐¥6 / 分享罐¥2\n6:3:1比例，AI给出参考建议"),
    (C_GREEN,  "✅ 劳动奖励机制",
     "义务家务不计薪（培养责任）/ 额外任务有报酬\n财商绘本讲给AI听 +¥10"),
    (RGBColor(0x8B, 0x5C, 0xF6), "🤖 AI财商学姐",
     "发零花钱时推送财商故事\n月度宝贝成长报告 / 18岁前完成基础财商体系"),
]
for i, (color, title, content) in enumerate(fn_cards):
    tp = 1.6 + i * 1.65
    add_rect(slide, left=4.0, top=tp, width=9.0, height=1.5,
             fill_color=C_WHITE, line_color=color, line_width=2)
    add_rect(slide, left=4.0, top=tp, width=0.1, height=1.5,
             fill_color=color)
    add_textbox(slide, title,
                left=4.2, top=tp+0.1, width=8.5, height=0.45,
                font_size=13, bold=True, color=color)
    add_textbox(slide, content,
                left=4.2, top=tp+0.55, width=8.5, height=0.85,
                font_size=11, color=C_BODY)

print("  Slide 11 完成")

# ═══════════════════════════════════════════════════════════════
# Slide 12 — 关键问题思考
# ═══════════════════════════════════════════════════════════════
print("生成 Slide 12 — 关键问题思考...")
slide = prs.slides.add_slide(blank_layout)
add_title(slide, "关键问题思考", slide_num=12)

qa_cards = [
    (C_RED,    RGBColor(0xFF, 0xFF, 0xFF), "01",
     "数据隐私边界：如何在本地加密与云端AI之间取得平衡？",
     "家庭财务数据极度敏感。芽计的答案：核心数据端侧AES-256加密，AI仅接触脱敏摘要，用户可选择断开云同步。"),
    (C_ORANGE, RGBColor(0xFF, 0xF9, 0xF5), "02",
     "AI建议的监管合规：财务建议是否构成「投资顾问」？",
     "中国监管对投资建议有严格规定。芽计定位为「决策辅助工具」而非投资顾问，内容经合规审查，核心输出均为信息参考。"),
    (C_GREEN,  RGBColor(0xF5, 0xFF, 0xF8), "03",
     "如何避免成为新的焦虑制造者？",
     "产品设计第一原则——先帮家长说'不'，再说'怎么做'。月度小报的第一行永远是'本月您帮孩子节省了¥X焦虑消费'。"),
]

for i, (border_color, bg_color, num, q, a) in enumerate(qa_cards):
    tp = 1.4 + i * 1.75
    add_rect(slide, left=0.7, top=tp, width=11.5, height=1.6,
             fill_color=bg_color, line_color=border_color, line_width=1)
    add_rect(slide, left=0.7, top=tp, width=0.1, height=1.6,
             fill_color=border_color)
    add_textbox(slide, num,
                left=0.9, top=tp+0.4, width=0.8, height=0.8,
                font_size=32, bold=True, color=border_color)
    add_textbox(slide, q,
                left=1.8, top=tp+0.1, width=10.2, height=0.45,
                font_size=13, bold=True, color=C_DARK_BLUE)
    add_textbox(slide, a,
                left=1.8, top=tp+0.6, width=10.2, height=0.85,
                font_size=11, color=C_BODY)

# 底部金句
add_rect(slide, left=0.7, top=6.55, width=11.5, height=0.55,
         fill_color=RGBColor(0xEE, 0xF3, 0xFF))
add_textbox(slide,
            '"我们相信，最好的财务顾问不是推销产品，而是帮你想清楚。" — 芽计 YaPlan',
            left=0.7, top=6.58, width=11.5, height=0.48,
            font_size=14, bold=True, color=C_DARK_BLUE, align='center')

print("  Slide 12 完成")

# ─────────────────────────────────────────
# 保存
# ─────────────────────────────────────────
print(f"\n正在保存到 {OUTPUT_PATH} ...")
prs.save(OUTPUT_PATH)
print("✅ PPT已生成")
